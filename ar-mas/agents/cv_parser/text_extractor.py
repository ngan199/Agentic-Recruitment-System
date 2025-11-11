"""
text_extractor.py

Provides:
- extract_pdf(path) -> dict
- extract_docx(path) -> dict
- extract_pptx(path) -> dict

Each returns a structured dict containing text blocks, tables, images metadata, and metadata
about extraction methods. All functions use shared helpers: normalize_text, ocr_image_bytes, etc.

Best-practice notes:
- OCR performed only when a page/slide is low-text or image-only to save CPU.
- PDF pages are rendered via PyMuPDF to images for page-level OCR fallback (no external renderer required).
- PPTX uses python-pptx for structure and falls back to image OCR for images or low-text slides.
- DOCX extracts headers, footers, body paragraphs, and tables.
"""

import io
import logging
import os
import re
import traceback
import unicodedata
from typing import Any, Dict, List, Tuple, Optional

from PIL import Image, ImageOps
import pytesseract

# External libraries
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document as DocxDocument

# ----------------------------
# Logging / Configuration
# ----------------------------
logger = logging.getLogger("text_extractor")
handler = logging.StreamHandler()
formatter = logging.Formatter("%(asctime)s %(levelname)s %(name)s - %(message)s")
handler.setFormatter(formatter)
logger.addHandler(handler)
logger.setLevel(logging.INFO)

# Configurable thresholds
MIN_TEXT_THRESHOLD = 200        # characters under which we consider OCR fallback
OCR_LANG = "eng"                # pytesseract language
MAX_IMAGE_PIXELS = 10_000_000   # reject extremely large images for OCR to avoid OOM
DPI_FOR_PDF_RENDER = 150        # resolution used to render PDF pages for OCR (ppi)
MAX_PDF_PAGES_FOR_OCR = 200     # safety cap
MAX_EXTRACTED_IMAGES = 200      # safety cap for images per file


def extract_file(path: str, prefer_text: bool = True, min_text_threshold: int = None) -> Dict[str, Any]:
    """
    Main dispatcher function to extract content from any supported document type.

    Args:
        path (str): Path to the input file (PDF, DOCX, PPTX).
        prefer_text (bool): If True, prefer textual content (over OCR fallback if available).
        min_text_threshold (int, optional): Minimum text length to trigger OCR fallback for PDFs/PPTX.

    Returns:
        Dict[str, Any]: Structured extraction result.
    
    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file type is unsupported.
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"File not found: {path}")

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        result = extract_pdf(path, min_text_threshold=min_text_threshold or 200)
    elif ext == ".docx":
        result = extract_docx(path)
    elif ext == ".pptx":
        result = extract_pptx(path, min_text_threshold=min_text_threshold or 200)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported types: .pdf, .docx, .pptx")

    # Optionally, if prefer_text is True, return a combined text string instead of full structured dict
    if prefer_text:
        if ext == ".pdf":
            from .text_extractor import extract_text_from_pdf_str
            return {"text": extract_text_from_pdf_str(path)}
        elif ext == ".docx":
            from .text_extractor import extract_text_from_docx_str
            return {"text": extract_text_from_docx_str(path)}
        elif ext == ".pptx":
            from .text_extractor import extract_text_from_pptx_str
            return {"text": extract_text_from_pptx_str(path)}

    return result


# ----------------------------
# Utilities
# ----------------------------
def normalize_text(text: Optional[str]) -> str:
    """Unicode-normalize, collapse whitespace, normalize newlines, and strip."""
    if not text:
        return ""
    # Unicode normalize to NFKC for compatibility
    text = unicodedata.normalize("NFKC", text)
    # normalize CRLF -> LF
    text = re.sub(r"\r\n?", "\n", text)
    # collapse repeated newlines > 2 into exactly two
    text = re.sub(r"\n{3,}", "\n\n", text)
    # collapse multiple spaces/tabs
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()


def safe_image_from_bytes(img_bytes: bytes) -> Optional[Image.Image]:
    """Return PIL.Image or None if invalid/too large."""
    try:
        with Image.open(io.BytesIO(img_bytes)) as im:
            # Basic check of pixel count to avoid OOM
            w, h = im.size
            if w * h > MAX_IMAGE_PIXELS:
                logger.warning(f"Image too large for OCR ({w}x{h}={w*h} pixels). Skipping OCR.")
                return None
            im = im.convert("RGB")
            return im.copy()
    except Exception as e:
        logger.debug(f"safe_image_from_bytes: invalid image - {e}")
        return None


def ocr_image_bytes(img_bytes: bytes, lang: str = OCR_LANG, preprocess: bool = True) -> str:
    """
    Run OCR on bytes using pytesseract and return normalized text.
    Includes optional preprocessing (grayscale & auto-contrast).
    """
    try:
        im = safe_image_from_bytes(img_bytes)
        if im is None:
            return ""
        if preprocess:
            # convert to grayscale & auto-contrast (simple but often helps)
            im = ImageOps.grayscale(im)
            im = ImageOps.autocontrast(im)
        text = pytesseract.image_to_string(im, lang=lang)
        return normalize_text(text)
    except Exception as e:
        logger.warning(f"OCR failed: {e}")
        return ""


# ----------------------------
# PDF Extraction
# ----------------------------
def extract_pdf(path: str, min_text_threshold: int = MIN_TEXT_THRESHOLD) -> Dict[str, Any]:
    """
    Extract text, tables (best-effort by text extraction), and images from a PDF using PyMuPDF.
    Performs OCR fallback per page by rendering page to PNG and OCRing when page text is too short.
    Returns a structured dict:
    {
        "filename": path,
        "pages": [
            {"page_number": n, "text": "...", "ocr_text": "...", "images": [ {meta} ], "extraction_method": "text|ocr" }
        ],
        "metadata": {...}
    }
    """
    result = {"filename": path, "pages": [], "metadata": {"num_pages": 0, "extraction_method": "pymupdf+ocr-fallback"}}
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    try:
        doc = fitz.open(path)
    except Exception as e:
        logger.exception("Failed to open PDF with PyMuPDF")
        raise e

    result["metadata"]["num_pages"] = len(doc)
    for pageno, page in enumerate(doc, start=1):
        page_entry: Dict[str, Any] = {"page_number": pageno, "text": "", "ocr_text": "", "images": [], "extraction_method": "text"}
        try:
            # Primary text extraction
            # 'text' option preserves some layout (lines)
            extracted_text = page.get_text("text") or ""
            extracted_text = normalize_text(extracted_text)
            page_entry["text"] = extracted_text

            # Extract images (embedded images)
            try:
                img_list = page.get_images(full=True) or []
                images_meta = []
                for img_index, img_info in enumerate(img_list):
                    # img_info format: (xref, smth, smth...)
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image.get("image")
                    if img_bytes:
                        images_meta.append({
                            "image_id": f"p{pageno}_img{img_index}",
                            "content_type": base_image.get("ext"),
                            "width": base_image.get("width"),
                            "height": base_image.get("height"),
                            "bytes": img_bytes,
                        })
                    if len(images_meta) >= MAX_EXTRACTED_IMAGES:
                        logger.warning("Reached MAX_EXTRACTED_IMAGES for PDF; skipping remaining images.")
                        break
                page_entry["images"] = images_meta
            except Exception as e:
                logger.debug(f"Failed to extract embedded images on page {pageno}: {e}")

            # Decide if we need OCR fallback:
            text_len = len(page_entry["text"].strip())
            if text_len < min_text_threshold:
                # Render page to pixmap (image) for OCR fallback
                try:
                    if len(doc) > MAX_PDF_PAGES_FOR_OCR:
                        # safety: avoid OCRing exceedingly large multi-page PDFs by default
                        logger.info(f"Skipping OCR rendering on page {pageno} because doc has {len(doc)} pages (above cap).")
                    else:
                        mat = fitz.Matrix(DPI_FOR_PDF_RENDER / 72.0, DPI_FOR_PDF_RENDER / 72.0)
                        pix = page.get_pixmap(matrix=mat, alpha=False)
                        img_bytes = pix.tobytes("png")
                        ocr_text = ocr_image_bytes(img_bytes)
                        if ocr_text:
                            page_entry["ocr_text"] = ocr_text
                            # If OCR produced more text than text extraction, mark as method
                            if len(ocr_text) > text_len:
                                page_entry["extraction_method"] = "ocr"
                        else:
                            logger.debug(f"OCR produced no text for page {pageno}.")
                except Exception as e:
                    logger.warning(f"Failed to render/ocr page {pageno}: {e}")

            # final normalization: prefer text extraction but keep ocr_text separate
            page_entry["text"] = normalize_text(page_entry.get("text", "") or "")
            page_entry["ocr_text"] = normalize_text(page_entry.get("ocr_text", "") or "")
        except Exception as e:
            logger.error(f"Error processing page {pageno}: {e}\n{traceback.format_exc()}")
            page_entry["error"] = str(e)
        result["pages"].append(page_entry)

    try:
        doc.close()
    except Exception:
        pass

    return result


# ----------------------------
# DOCX Extraction
# ----------------------------
def extract_docx(path: str, min_text_threshold: int = MIN_TEXT_THRESHOLD) -> Dict[str, Any]:
    """
    Extract text and tables from a .docx file. Returns structure:
    {
       "filename": path,
       "paragraphs": [...],
       "tables": [ { "rows": [[cell_text,...], ...] } ],
       "headers": ...,
       "footers": ...,
       "metadata": {...}
    }
    """
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    result: Dict[str, Any] = {
        "filename": path,
        "paragraphs": [],
        "tables": [],
        "headers": [],
        "footers": [],
        "metadata": {"extraction_method": "python-docx"},
    }

    try:
        doc = DocxDocument(path)
    except Exception as e:
        logger.exception("Failed to open DOCX")
        raise e

    # Extract body paragraphs (keeping simple order)
    try:
        paragraphs = []
        for para in doc.paragraphs:
            t = normalize_text(para.text)
            if t:
                paragraphs.append(t)
        result["paragraphs"] = paragraphs
    except Exception as e:
        logger.debug(f"Failed to read paragraphs: {e}")

    # Extract tables
    try:
        tables = []
        for tbl in doc.tables:
            rows = []
            for r in tbl.rows:
                row_cells = []
                for cell in r.cells:
                    # cell.text already concatenates paragraphs inside
                    row_cells.append(normalize_text(cell.text))
                rows.append(row_cells)
            tables.append({"rows": rows})
        result["tables"] = tables
    except Exception as e:
        logger.debug(f"Failed to extract tables from DOCX: {e}")

    # Headers & footers (best-effort; python-docx supports sections' header/footer)
    try:
        for section in doc.sections:
            header = []
            if section.header:
                for p in section.header.paragraphs:
                    if p.text:
                        header.append(normalize_text(p.text))
            result["headers"].append("\n".join(header))

            footer = []
            if section.footer:
                for p in section.footer.paragraphs:
                    if p.text:
                        footer.append(normalize_text(p.text))
            result["footers"].append("\n".join(footer))
    except Exception as e:
        logger.debug(f"Failed to read headers/footers: {e}")

    return result


# ----------------------------
# PPTX Extraction (reuses OCR helpers)
# ----------------------------
def _pptx_shape_has_text(shape) -> bool:
    return getattr(shape, "has_text_frame", False)


def _pptx_extract_text_from_shape(shape) -> str:
    """Extract paragraph/run text from a shape."""
    try:
        if not _pptx_shape_has_text(shape):
            return ""
        lines = []
        for para in shape.text_frame.paragraphs:
            runs = [r.text for r in para.runs if r.text]
            line = "".join(runs) if runs else para.text
            if line:
                lines.append(line)
        return "\n".join(lines)
    except Exception as e:
        logger.debug(f"_pptx_extract_text_from_shape error: {e}")
        return ""


def _pptx_extract_table(shape) -> List[List[str]]:
    table_rows = []
    try:
        tbl = shape.table
        for r in tbl.rows:
            row = [normalize_text(cell.text) for cell in r.cells]
            table_rows.append(row)
    except Exception as e:
        logger.debug(f"_pptx_extract_table error: {e}")
    return table_rows


def _pptx_extract_image(shape, slide_idx: int, img_idx: int) -> Dict[str, Any]:
    """Return image metadata (bytes + dims) from a picture shape."""
    try:
        image = shape.image
        blob = image.blob
        return {
            "type": "image",
            "image_id": f"slide{slide_idx}_img{img_idx}",
            "ext": image.ext,
            "content_type": image.content_type,
            "left": getattr(shape, "left", None),
            "top": getattr(shape, "top", None),
            "width": getattr(shape, "width", None),
            "height": getattr(shape, "height", None),
            "bytes": blob,
        }
    except Exception as e:
        logger.debug(f"_pptx_extract_image error: {e}")
        return {}


def _pptx_shape_to_blocks(shape, slide_idx: int, img_counter: List[int]) -> List[Dict[str, Any]]:
    """Recursively convert a shape to one or more blocks (text/table/image)."""
    blocks: List[Dict[str, Any]] = []
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                blocks.extend(_pptx_shape_to_blocks(s, slide_idx, img_counter))
            return blocks

        if _pptx_shape_has_text(shape):
            text = normalize_text(_pptx_extract_text_from_shape(shape))
            if text:
                blocks.append({
                    "type": "text",
                    "text": text,
                    "coords": {"left": getattr(shape, "left", 0), "top": getattr(shape, "top", 0)},
                    "source": "shape",
                })
            return blocks

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_rows = _pptx_extract_table(shape)
            if table_rows:
                blocks.append({
                    "type": "table",
                    "rows": table_rows,
                    "coords": {"left": getattr(shape, "left", 0), "top": getattr(shape, "top", 0)},
                    "source": "table",
                })
            return blocks

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_meta = _pptx_extract_image(shape, slide_idx, img_counter[0])
            if img_meta:
                blocks.append(img_meta)
                img_counter[0] += 1
            return blocks

        return blocks
    except Exception as e:
        logger.debug(f"_pptx_shape_to_blocks error: {e}")
        return blocks


def extract_pptx(path: str, min_text_threshold: int = MIN_TEXT_THRESHOLD, include_hidden: bool = True) -> Dict[str, Any]:
    """
    Extract structured content from PPTX:
    returns:
    {
      "filename": path,
      "slides": [
         {"index": i, "title": "...", "blocks": [ {type:text/table/image/ocr_text}, ... ], "notes": "...", "extraction_method": "..."}
      ],
      "metadata": {...}
    }
    """
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    result: Dict[str, Any] = {"filename": path, "slides": [], "metadata": {"extraction_method": "python-pptx+ocr-fallback"}}
    try:
        prs = Presentation(path)
    except Exception as e:
        logger.exception("Failed to open PPTX")
        raise e

    for i, slide in enumerate(prs.slides):
        slide_entry: Dict[str, Any] = {"index": i, "title": "", "blocks": [], "notes": "", "extraction_method": "text"}
        try:
            # skip hidden slides if desired (python-pptx doesn't expose .hidden property directly, so skipping is optional)
            img_counter = [0]
            for shape in slide.shapes:
                slide_entry["blocks"].extend(_pptx_shape_to_blocks(shape, i, img_counter))

            # Notes
            try:
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_entry["notes"] = normalize_text(slide.notes_slide.notes_text_frame.text)
            except Exception:
                slide_entry["notes"] = ""

            # Determine title if present (placeholder with idx 0 often)
            try:
                # best-effort: find the top-most text or placeholder typical of title
                title_candidates = [b for b in slide_entry["blocks"] if b.get("type") == "text"]
                if title_candidates:
                    # choose the one with smallest top coord as title heuristic
                    title_candidates_sorted = sorted(title_candidates, key=lambda x: x.get("coords", {}).get("top", 0))
                    slide_entry["title"] = title_candidates_sorted[0]["text"][:300]
            except Exception:
                pass

            # If total text is tiny and images exist, OCR each image + optionally append ocr_text block
            total_text_len = sum(len(b.get("text", "")) for b in slide_entry["blocks"] if b.get("type") == "text")
            if total_text_len < min_text_threshold:
                ocr_texts = []
                for b in slide_entry["blocks"]:
                    if b.get("type") == "image" and b.get("bytes"):
                        ocr_t = ocr_image_bytes(b["bytes"])
                        if ocr_t:
                            b["ocr_text"] = ocr_t
                            ocr_texts.append(ocr_t)
                if ocr_texts:
                    slide_entry["extraction_method"] = "ocr"
                    slide_entry["blocks"].append({
                        "type": "ocr_text",
                        "text": normalize_text("\n".join(ocr_texts)),
                        "source": "ocr_fallback"
                    })

            # Final normalization / ordering by (top,left)
            for b in slide_entry["blocks"]:
                if b.get("type") == "text":
                    b["text"] = normalize_text(b.get("text", ""))
            slide_entry["blocks"].sort(key=lambda b: (b.get("coords", {}).get("top", 0), b.get("coords", {}).get("left", 0)))
        except Exception as e:
            slide_entry["error"] = str(e)
            logger.error(f"PPTX slide {i} processing error: {e}\n{traceback.format_exc()}")

        result["slides"].append(slide_entry)
    return result


# ----------------------------
# Convenience wrappers
# ----------------------------
def extract_text_from_pdf_str(path: str) -> str:
    """Compatibility wrapper returning combined text string (prefers text extraction, falls back to OCR pages)."""
    doc = extract_pdf(path)
    parts = []
    for p in doc["pages"]:
        if p.get("text"):
            parts.append(p["text"])
        elif p.get("ocr_text"):
            parts.append(p["ocr_text"])
    return normalize_text("\n\n".join(parts))


def extract_text_from_pptx_str(path: str) -> str:
    """Return combined text of PPTX (text + notes + OCR fallbacks)."""
    doc = extract_pptx(path)
    parts = []
    for s in doc["slides"]:
        for b in s.get("blocks", []):
            if b.get("type") in ("text", "ocr_text"):
                parts.append(b.get("text", ""))
        if s.get("notes"):
            parts.append(s["notes"])
    return normalize_text("\n\n".join(parts))


def extract_text_from_docx_str(path: str) -> str:
    doc = extract_docx(path)
    parts = []
    parts.extend(doc.get("headers", []))
    parts.extend(doc.get("paragraphs", []))
    for t in doc.get("tables", []):
        # flatten tables to text rows
        for r in t.get("rows", []):
            parts.append(" | ".join(r))
    parts.extend(doc.get("footers", []))
    return normalize_text("\n\n".join([p for p in parts if p]))

